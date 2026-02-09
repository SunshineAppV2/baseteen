from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_presentation():
    prs = Presentation()

    # Image Paths (Hardcoded based on generation)
    # Using the exact filenames from the generation steps
    image_paths = {
        0: r"C:\Users\USUARIO\.gemini\antigravity\brain\e07fe714-b837-4ff5-801c-5867babb5f27\min_slide_1_cover_1769953112508.png",
        1: r"C:\Users\USUARIO\.gemini\antigravity\brain\e07fe714-b837-4ff5-801c-5867babb5f27\min_slide_2_challenge_1769953125963.png",
        2: r"C:\Users\USUARIO\.gemini\antigravity\brain\e07fe714-b837-4ff5-801c-5867babb5f27\min_slide_3_solution_1769953139779.png",
        3: r"C:\Users\USUARIO\.gemini\antigravity\brain\e07fe714-b837-4ff5-801c-5867babb5f27\slide_4_management_1769952881108.png",
        4: r"C:\Users\USUARIO\.gemini\antigravity\brain\e07fe714-b837-4ff5-801c-5867babb5f27\slide_5_teacher_v2.png", 
        5: r"C:\Users\USUARIO\.gemini\antigravity\brain\e07fe714-b837-4ff5-801c-5867babb5f27\slide_6_motivation_v2.png", 
        6: r"C:\Users\USUARIO\.gemini\antigravity\brain\e07fe714-b837-4ff5-801c-5867babb5f27\slide_7_gamification_v2.png", 
        7: r"C:\Users\USUARIO\.gemini\antigravity\brain\e07fe714-b837-4ff5-801c-5867babb5f27\slide_8_ranking_1769952903626.png",
        8: r"C:\Users\USUARIO\.gemini\antigravity\brain\e07fe714-b837-4ff5-801c-5867babb5f27\slide_9_ux_1769952920161.png",
        9: r"C:\Users\USUARIO\.gemini\antigravity\brain\e07fe714-b837-4ff5-801c-5867babb5f27\slide_10_coordinator_1769952934422.png",
        10: r"C:\Users\USUARIO\.gemini\antigravity\brain\e07fe714-b837-4ff5-801c-5867babb5f27\slide_11_security_1769952953480.png",
        11: r"C:\Users\USUARIO\.gemini\antigravity\brain\e07fe714-b837-4ff5-801c-5867babb5f27\slide_12_next_steps_1769952979078.png"
    }

    # Define slides data
    slides_data = [
        {
            "title": "BaseTeen: A Revolução na Gestão da Sua Base",
            "content": "Simplifique a administração e potencialize o engajamento dos seus adolescentes.",
            "type": "title"
        },
        {
            "title": "O Desafio dos Coordenadores",
            "content": 
                "• Dificuldade no acompanhamento manual dos alunos.\n"
                "• Baixo engajamento no estudo diário da lição.\n"
                "• Falta de métricas claras sobre a saúde da base.\n"
                "• Processos administrativos lentos e burocráticos.",
            "type": "bullet"
        },
        {
            "title": "A Solução BaseTeen",
            "content":
                "Uma plataforma integrada que une Gestão Eficiente e Gamificação.\n\n"
                "• Painel Administrativo intuitivo.\n"
                "• App para o adolescente.\n"
                "• Controle total na palma da mão.",
            "type": "bullet"
        },
        {
            "title": "Facilidade de Gerenciamento",
            "content":
                "Diga adeus às planilhas complexas.\n\n"
                "• Cadastro Simplificado: Links de auto-cadastro.\n"
                "• Controle de Assinaturas: Gestão automática.\n"
                "• Hierarquia Clara: Visualize Região, Distrito, Base.",
            "type": "bullet"
        },
        {
            "title": "Gestão de Professores de Classe",
            "content":
                "Empodere seus professores com dados reais.\n\n"
                "• Chamada Inteligente: Registre a presença em segundos.\n"
                "• Termômetro de Crescimento: Meça a evolução da classe pela frequência.\n"
                "• Relatórios automáticos de evasão e retenção.",
            "type": "bullet"
        },
        {
            "title": "Motivação para o Estudo Diário",
            "content":
                "Transforme o dever em desejo.\n\n"
                "• Incentivo Diário: Notificações e metas.\n"
                "• Recompensas Imediatas: XP e badges.\n"
                "• O estudo deixa de ser 'chato' e vira conquista.",
            "type": "bullet"
        },
        {
            "title": "Gamificação e Engajamento",
            "content":
                "A linguagem que o adolescente entende.\n\n"
                "• Sistema de Níveis: Evolução constante.\n"
                "• Missões e Desafios: Tarefas extras.\n"
                "• Engajamento orgânico através da diversão.",
            "type": "bullet"
        },
        {
            "title": "O Poder do Ranking",
            "content":
                "Competição saudável que gera crescimento.\n\n"
                "• Ranking em Tempo Real: Atualizado a cada ação.\n"
                "• Visualização Global e Local: Compare desempenhos.\n"
                "• Destaque para os alunos mais dedicados.",
            "type": "bullet"
        },
        {
            "title": "Usabilidade Incrível (UX)",
            "content":
                "Feito para ser amado pelos usuários.\n\n"
                "• Design Premium: Moderno e responsivo.\n"
                "• Mobile-First: Perfeito em qualquer celular.\n"
                "• Navegação fluida que não trava a rotina.",
            "type": "bullet"
        },
        {
            "title": "Para Você, Coordenador",
            "content":
                "Mais tempo para as pessoas, menos para o papel.\n\n"
                "• Métricas de Crescimento: Gráficos baseados nas chamadas.\n"
                "• Tomada de decisão baseada em dados reais.\n"
                "• Identifique tendências de queda antes que aconteçam.",
            "type": "bullet"
        },
        {
            "title": "Tecnologia e Segurança",
            "content":
                "Confiabilidade total.\n\n"
                "• Dados Seguros: Proteção LGPD.\n"
                "• Sempre Disponível: Sistema em nuvem.\n"
                "• Suporte técnico dedicado.",
            "type": "bullet"
        },
        {
            "title": "Próximos Passos",
            "content":
                "Vamos começar a transformação?\n\n"
                "1. Faça seu cadastro hoje.\n"
                "2. Cadastre sua equipe de professores.\n"
                "3. Convide seus alunos.\n\n"
                "BaseTeen - Conectando Gerações.",
            "type": "bullet"
        }
    ]

    for i, slide_info in enumerate(slides_data):
        if slide_info["type"] == "title":
            slide_layout = prs.slide_layouts[1] # Using Title and Content for consistency with images
        else:
            slide_layout = prs.slide_layouts[1]

        slide = prs.slides.add_slide(slide_layout)
        
        # Set Title
        title = slide.shapes.title
        title.text = slide_info["title"]

        # Set Content
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear() # Clear default bullets
        tf.word_wrap = True
        
        # Add content text (adjusting size to fit image)
        p = tf.add_paragraph()
        p.text = slide_info["content"]
        p.font.size = Pt(18)

        # Add Image
        if i in image_paths:
            img_path = image_paths[i]
            
            # Check if file exists to avoid crashes (especially for the retried ones)
            # For the retried ones, we might need to find the exact filename if it has a timestamp
            if not os.path.exists(img_path):
                # Try to find a file starting with the prefix in the directory
                dir_name = os.path.dirname(img_path)
                base_name = os.path.basename(img_path).split('.')[0] # Remove extension
                prefix = base_name.split('_v2')[0] if '_v2' in base_name else base_name
                
                found = False
                if os.path.exists(dir_name):
                    for f in os.listdir(dir_name):
                        if f.startswith(prefix) and f.endswith(".png"):
                           img_path = os.path.join(dir_name, f)
                           found = True
                           break
                if not found:
                    print(f"Warning: Image for slide {i} not found: {img_path}")
                    continue

            # Position image on the right side
            left = Inches(5.5)
            top = Inches(2)
            width = Inches(4)
            # height = Inches(3) # Let aspect ratio decide
            
            try:
                slide.shapes.add_picture(img_path, left, top, width=width)
            except Exception as e:
                print(f"Error adding image {img_path}: {e}")

    output_path = "presentation_baseteen.pptx"
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    create_presentation()
